#===============================================================================
#
#  Flatmap viewer and annotation tools
#
#  Copyright (c) 2018 - 2023  David Brooks
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
#===============================================================================

from typing import Optional

#===============================================================================

from lxml import etree
import numpy as np
import shapely.geometry
import shapely.ops

from pptx import Presentation
from pptx.dml.fill import FillFormat
from pptx.enum.dml import MSO_FILL_TYPE             # type: ignore
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN     # type: ignore
from pptx.shapes.autoshape import Shape as PptxShape
from pptx.shapes.connector import Connector as PptxConnector
from pptx.shapes.group import GroupShape as PptxGroupShape
from pptx.shapes.shapetree import GroupShapes as PptxGroupShapes
from pptx.shapes.shapetree import SlideShapes as PptxSlideShapes
from pptx.slide import Slide as PptxSlide

#===============================================================================

from mapmaker.annotation import Annotator
from mapmaker.geometry import Transform
from mapmaker.properties.markup import parse_layer_directive, parse_markup
from mapmaker.sources import MapBounds, WORLD_METRES_PER_EMU
from mapmaker.sources.shape import Shape, SHAPE_TYPE
from mapmaker.sources.shapefilter import ShapeFilter
from mapmaker.utils import FilePath, log, ProgressBar, TreeList

from .colour import ColourMap, ColourTheme
from .geometry import get_shape_geometry
from .presets import DRAWINGML, PPTX_NAMESPACE, pptx_resolve, pptx_uri
from .transform import DrawMLTransform

#===============================================================================

class PowerpointShape(Shape):
    pass

#===============================================================================

# (colour, opacity)
ColourPair = tuple[Optional[str], float]

#===============================================================================

class Slide:
    def __init__(self, source_id: str, kind: str, index: int, pptx_slide: PptxSlide,
                 theme: ColourTheme, bounds: MapBounds, transform: Transform,
                 shape_filter: Optional[ShapeFilter]=None):
        self.__source_id = source_id
        self.__kind = kind
        self.__id = 'slide-{:02d}'.format(index+1)
        # Get any layer directives
        if pptx_slide.has_notes_slide:
            notes_slide = pptx_slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text
            if notes_text.startswith('.'):
                layer_directive = parse_layer_directive(notes_text)
                if 'error' in layer_directive:
                    log.error('error', f'Slide {index+1}: invalid layer directive: {notes_text}')
                if 'id' in layer_directive:
                    self.__id = layer_directive['id']
        self.__colour_map = ColourMap(theme, pptx_slide)
        self.__pptx_slide = pptx_slide
        self.__geometry = shapely.geometry.box(*bounds)
        self.__transform = transform
        self.__shapes = TreeList()
        self.__shapes_by_id: dict[str, PowerpointShape] = {}

    @property
    def colour_map(self) -> ColourMap:
        return self.__colour_map

    @property
    def geometry(self) -> shapely.geometry.base.BaseGeometry:
        return self.__geometry

    @property
    def kind(self) -> str:
        return self.__kind

    @property
    def id(self) -> str:
        return self.__id

    @property
    def pptx_slide(self) -> PptxSlide:
        return self.__pptx_slide

    @property
    def shapes(self) -> TreeList:
        return self.__shapes

    @property
    def slide_id(self) -> int:
        return self.__pptx_slide.slide_id

    @property
    def source_id(self):
        return self.__source_id

    def shape(self, id: str) -> Optional[PowerpointShape]:
    #=====================================================
        return self.__shapes_by_id.get(id)

    def __new_shape(self, type, id: int, geometry, properties=None) -> PowerpointShape:
    #==================================================================================
        shape_id = f'{self.__id}/{id}'
        shape = (PowerpointShape(type, shape_id, geometry, properties) if properties is not None
            else PowerpointShape(type, shape_id, geometry))
        self.__shapes_by_id[shape_id] = shape
        return shape

    def process(self, annotator: Annotator) -> TreeList:
    #===================================================
        # Return the slide's group structure as a nested list of Shapes
        self.__shapes = TreeList([self.__new_shape(SHAPE_TYPE.GROUP, 'root', self.__geometry)])
        self.__shapes.extend(self.__process_pptx_shapes(self.__pptx_slide.shapes,      # type: ignore
                                                        self.__transform, show_progress=True))
        return self.__shapes

    def __get_colour(self, shape: PptxConnector | PptxGroupShape | PptxShape,
                     group_colour: Optional[ColourPair]=None) -> ColourPair:
    #=======================================================================
        def colour_from_fill(shape, fill) -> ColourPair:
            if fill.type == MSO_FILL_TYPE.SOLID:                    # type: ignore
                return (self.__colour_map.lookup(fill.fore_color),
                        fill.fore_color.alpha)
            elif fill.type == MSO_FILL_TYPE.GRADIENT:               # type: ignore
                log.warning(f'{shape.name}: gradient fill ignored')
            elif fill.type == MSO_FILL_TYPE.GROUP:                  # type: ignore
                if group_colour is not None:
                    return group_colour
            elif fill.type is not None and fill.type != MSO_FILL_TYPE.BACKGROUND:   # type: ignore
                log.warning(f'{shape.name}: unsupported fill type: {fill.type}')
            return (None, 1.0)

        colour = None
        alpha = 1.0
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:                # type: ignore
            colour, alpha = colour_from_fill(shape, FillFormat.from_fill_parent(shape.element.grpSpPr))
        elif shape.shape_type != MSO_SHAPE_TYPE.LINE:               # type: ignore
            colour, alpha = colour_from_fill(shape, shape.fill)     # type: ignore
        elif shape.line.fill.type == MSO_FILL_TYPE.SOLID:           # type: ignore
            colour = self.__colour_map.lookup(shape.line.color)     # type: ignore
            alpha = shape.line.fill.fore_color.alpha                # type: ignore
        elif shape.line.fill.type is None:                          # type: ignore
            # Check for a fill colour in the <style> block
            xml = etree.fromstring(shape.element.xml)
            if (scheme_colour := xml.find('.//p:style/a:fillRef/a:schemeClr',
                                            namespaces=PPTX_NAMESPACE)) is not None:
                colour = self.__colour_map.scheme_colour(scheme_colour.attrib['val'])
        elif shape.line.fill.type != MSO_FILL_TYPE.BACKGROUND:      # type: ignore
            log.warning(f'{shape.name}: unsupported line fill type: {shape.line.fill.type}')    # type: ignore
        return (colour, alpha)


    def __shapes_as_group(self, group: PptxGroupShape, shapes: TreeList) -> Shape | TreeList:
    #========================================================================================
        if len(shapes) < 2:  ## shapes[0] might be a TreeList ##
            return shapes    ## or shapes[0].type != SHAPE_TYPE.FEATURE:
        colour = shapes[0].colour
        label = shapes[0].label
        alignment = shapes[0].properties.get('align')
        geometry = [shapes[0].geometry]
        for shape in shapes[1:]:
            if (isinstance(shape, TreeList)
             or shape.type != SHAPE_TYPE.FEATURE
             or colour != shape.colour):
                return shapes
            if label == '':
                label = shape.label
                alignment = shape.properties.get('align')
            elif shape.label != '' and label != shape.label:
                return shapes
            geometry.append(shape.geometry)
        # Merge a group of shapes that are all the same colour
        # having a common label into a single shape
        return Shape(SHAPE_TYPE.FEATURE, group.shape_id,
                      shapely.ops.unary_union(geometry), {
                        'colour': colour,
                        'label': label,
                        'shape-name': group.name,
                        'text-align': alignment
                       })

    def __process_group(self, group: PptxGroupShape, transform: Transform) -> Shape | TreeList:
    #==========================================================================================
        colour = self.__get_colour(group)
        group_shapes = self.__shapes_as_group(group,
                            self.__process_pptx_shapes(group.shapes,        # type: ignore
                                transform@DrawMLTransform(group),
                                group_colour=colour))
        if isinstance(group_shapes, Shape):
            return group_shapes
        shapes = TreeList([self.__new_shape(SHAPE_TYPE.GROUP, group.shape_id, None, {
            'colour': colour[0],
            'opacity': colour[1],
            'pptx-shape': group
        })])
        shapes.extend(group_shapes)
        return shapes

    def __process_pptx_shapes(self, pptx_shapes: PptxGroupShapes | PptxSlideShapes,
                              transform: Transform, group_colour: Optional[ColourPair]=None,
                              show_progress=False) -> TreeList:
    #========================================================================================
        def text_alignment(shape) -> tuple[str, str]:
            para = shape.text_frame.paragraphs[0].alignment
            vertical = shape.text_frame.vertical_anchor
            return ('left' if para in [PP_ALIGN.LEFT, PP_ALIGN.DISTRIBUTE, PP_ALIGN.JUSTIFY, PP_ALIGN.JUSTIFY_LOW] else
                    'right' if para == PP_ALIGN.RIGHT else
                    'centre',
                    'top' if vertical == MSO_ANCHOR.TOP else
                    'bottom' if vertical == MSO_ANCHOR.BOTTOM else
                    'middle')
        def text_content(shape) -> str:
            text = shape.text.replace('\n', ' ').replace('\xA0', ' ').replace('\v', ' ').strip() # Newline, non-breaking space, vertical-tab
            return text if text not in ['', '.'] else ''

        progress_bar = ProgressBar(show=show_progress,
            total=len(pptx_shapes),
            unit='shp', ncols=40,
            bar_format='{l_bar}{bar}| {n_fmt}/{total_fmt}')
        shapes = TreeList()
        for pptx_shape in pptx_shapes:
            shape_name = pptx_shape.name
            shape_properties = parse_markup(shape_name) if shape_name.startswith('.') else {}
            if (pptx_shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE              # type: ignore
             or pptx_shape.shape_type == MSO_SHAPE_TYPE.FREEFORM                # type: ignore
             or pptx_shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX                # type: ignore
             or pptx_shape.shape_type == MSO_SHAPE_TYPE.LINE):                  # type: ignore
                colour, alpha = self.__get_colour(pptx_shape, group_colour)     # type: ignore
                shape_properties.update({
                    'shape-name': shape_name,
                    'colour': colour
                })
                if alpha < 1.0:
                    shape_properties['opacity'] = round(100*alpha, 1)
                geometry = get_shape_geometry(pptx_shape, transform, shape_properties)
                if geometry is not None and geometry.is_valid:
                    shape_xml = etree.fromstring(pptx_shape.element.xml)
                    for link_ref in shape_xml.findall('.//a:hlinkClick',
                                                    namespaces=PPTX_NAMESPACE):
                        r_id = link_ref.attrib[pptx_resolve('r:id')]
                        if (r_id in pptx_shape.part.rels
                         and pptx_shape.part.rels[r_id].reltype == pptx_uri('r:hyperlink')):
                            shape_properties['hyperlink'] = pptx_shape.part.rels[r_id].target_ref
                            break
                    if pptx_shape.shape_type == MSO_SHAPE_TYPE.LINE:            # type: ignore
                        ## cf. pptx2svg for stroke colour
                        shape_type = SHAPE_TYPE.CONNECTOR
                        if (connection := shape_xml.find('.//p:nvCxnSpPr/p:cNvCxnSpPr',
                                                        namespaces=PPTX_NAMESPACE)) is not None:
                            for c in connection.getchildren():
                                if c.tag == DRAWINGML('stCxn'):
                                    shape_properties['connection-start'] = int(c.attrib['id'])
                                elif c.tag == DRAWINGML('endCxn'):
                                    shape_properties['connection-end'] = int(c.attrib['id'])
                        shape_properties['line-style'] = pptx_shape.line.prstDash                   # type: ignore
                        shape_properties['head-end'] = pptx_shape.line.headEnd.get('type', 'none')  # type: ignore
                        shape_properties['tail-end'] = pptx_shape.line.tailEnd.get('type', 'none')  # type: ignore
                        shape_properties['stroke-width'] = abs(transform.scale_length((int(pptx_shape.line.width.emu), 0))[0])  # type: ignore
                    else:
                        shape_type = SHAPE_TYPE.FEATURE
                        label = text_content(pptx_shape)
                        if label != '':
                            shape_properties['label'] = label
                            shape_properties['align'] = text_alignment(pptx_shape)
                    shape_properties['pptx-shape'] = pptx_shape
                    shape = self.__new_shape(shape_type, pptx_shape.shape_id, geometry, shape_properties)
                    shapes.append(shape)
                elif geometry is None:
                    log.warning(f'Shape "{shape_name}" {pptx_shape.shape_type}/{shape_properties.get("shape-kind")} not processed -- cannot get geometry')
                else:
                    log.warning(f'Shape "{shape_name}" {pptx_shape.shape_type}/{shape_properties.get("shape-kind")} not processed -- cannot get valid geometry')
            elif pptx_shape.shape_type == MSO_SHAPE_TYPE.GROUP:             # type: ignore
                shapes.append(self.__process_group(pptx_shape, transform))  # type: ignore
            elif pptx_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:           # type: ignore
                log.warning('Image "{}" {} not processed...'.format(shape_name, str(pptx_shape.shape_type)))
            else:
                log.warning('Shape "{}" {} not processed...'.format(shape_name, str(pptx_shape.shape_type)))
            progress_bar.update(1)

        progress_bar.close()
        return shapes

#===============================================================================

class Powerpoint():
    def __init__(self, source_id: str, source_href: str, source_kind: str, shape_filter: Optional[ShapeFilter]=None,
                 SlideClass=Slide):
        ppt_bytes = FilePath(source_href).get_BytesIO()
        pptx = Presentation(ppt_bytes)

        (width, height) = (pptx.slide_width, pptx.slide_height)
        self.__transform = Transform([[WORLD_METRES_PER_EMU,                     0, 0],
                                      [                    0, -WORLD_METRES_PER_EMU, 0],
                                      [                    0,                     0, 1]])@np.array([[1.0, 0.0,  -width/2.0],
                                                                                                    [0.0, 1.0, -height/2.0],
                                                                                                    [0.0, 0.0,         1.0]])
        top_left = self.__transform.transform_point((0, 0))
        bottom_right = self.__transform.transform_point((width, height))
        # southwest and northeast corners
        self.__bounds = (top_left[0], bottom_right[1], bottom_right[0], top_left[1])

        colour_theme = ColourTheme(ppt_bytes)
        self.__slides: list[Slide] = [SlideClass(source_id,
                                                 source_kind,
                                                 slide_index,
                                                 slide,
                                                 colour_theme,
                                                 self.__bounds,
                                                 self.__transform,
                                                 shape_filter=shape_filter)
                                            for slide_index, slide in enumerate(pptx.slides)]

    @property
    def bounds(self) -> MapBounds:
        return self.__bounds

    @property
    def slides(self) -> list[Slide]:
        return self.__slides

    @property
    def transform(self) -> Transform:
        return self.__transform

#===============================================================================
