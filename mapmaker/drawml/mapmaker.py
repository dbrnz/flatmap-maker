#===============================================================================
#
#  Flatmap viewer and annotation tools
#
#  Copyright (c) 2019  David Brooks
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
#===============================================================================

from collections import defaultdict
from math import sqrt, sin, cos, pi as PI
import os

#===============================================================================

import numpy as np

import pptx.shapes.connector
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from tqdm import tqdm

#===============================================================================

from flatmap import MapLayer
from parser import Parser
from properties import JsonProperties

#===============================================================================

# Internal PPT units are EMUs (English Metric Units)

EMU_PER_CM  = 360000
EMU_PER_IN  = 914400

DOTS_PER_IN = 96

EMU_PER_DOT = EMU_PER_IN/DOTS_PER_IN

#===============================================================================

def cm_coords(x, y):
#===================
    return (x/EMU_PER_CM, y/EMU_PER_CM)

def ellipse_point(a, b, theta):
#==============================
    a_sin_theta = a*sin(theta)
    b_cos_theta = b*cos(theta)
    circle_radius = sqrt(a_sin_theta**2 + b_cos_theta**2)
    return (a*b_cos_theta/circle_radius, b*a_sin_theta/circle_radius)

#===============================================================================

class Transform(object):
    def __init__(self, shape, bbox=None):
        xfrm = shape.element.xfrm

        # From Section L.4.7.6 of ECMA-376 Part 1
        (Bx, By) = ((xfrm.chOff.x, xfrm.chOff.y)
                        if xfrm.chOff is not None else
                    (0, 0))
        (Dx, Dy) = ((xfrm.chExt.cx, xfrm.chExt.cy)
                        if xfrm.chExt is not None else
                    bbox)
        (Bx_, By_) = (xfrm.off.x, xfrm.off.y)
        (Dx_, Dy_) = (xfrm.ext.cx, xfrm.ext.cy)
        theta = xfrm.rot*PI/180.0
        Fx = -1 if xfrm.flipH else 1
        Fy = -1 if xfrm.flipV else 1
        T_st = np.array([[Dx_/Dx,      0, Bx_ - (Dx_/Dx)*Bx] if Dx != 0 else [1, 0, Bx_],
                         [     0, Dy_/Dy, By_ - (Dy_/Dy)*By] if Dy != 0 else [0, 1, By_],
                         [     0,      0,                 1]])
        U = np.array([[1, 0, -(Bx_ + Dx_/2.0)],
                      [0, 1, -(By_ + Dy_/2.0)],
                      [0, 0,                1]])
        R = np.array([[cos(theta), -sin(theta), 0],
                      [sin(theta),  cos(theta), 0],
                      [0,                    0, 1]])
        Flip = np.array([[Fx,  0, 0],
                         [ 0, Fy, 0],
                         [ 0,  0, 1]])
        T_rf = np.linalg.inv(U)@R@Flip@U
        self.__T = T_rf@T_st

    def matrix(self):
        return self.__T

#===============================================================================

class Feature(object):
    def __init__(self, feature_id, geometry, properties, has_children=False):
        self.__feature__id = feature_id
        self.__geometry = geometry
        self.__properties = properties.copy()
        self.__properties['featureId'] = feature_id   # Used by flatmap viewer
        self.__has_children = has_children

    def __str__(self):
        return 'Feature {}: {}'.format(self.__geometry.geom_type, self.__properties)

    @property
    def annotated(self):
        return self.shape_name.startswith('.')

    @property
    def feature_id(self):
        return self.__feature__id

    @property
    def geom_type(self):
        return self.__geometry.geom_type if self.__geometry else None

    @property
    def geometry(self):
        return self.__geometry

    @geometry.setter
    def geometry(self, geometry):
        self.__geometry = geometry

    @property
    def has_children(self):
        return self.__has_children

    @property
    def id(self):
        return self.__properties.get('id')

    @property
    def properties(self):
        return self.__properties

    @property
    def shape_id(self):
        return self.__feature__id.split('#')[-1]

    @property
    def shape_name(self):
        return self.__properties.get('shape_name', '')

    def has_property(self, property):
        return self.__properties.get(property, '') != ''

    def property(self, property, default=None):
        return self.__properties.get(property, default)

#===============================================================================

class FeaturesValueError(ValueError):
    def __init__(self, msg, features):
        super().__init__('\n  '.join([msg] + [str(f) for f in features]))

#===============================================================================

class SlideLayer(MapLayer):
    def __init__(self, mapmaker, slide, slide_number):
        super().__init__(slide_number, mapmaker)
        self.__slide = slide
        self.__slide_number = slide_number
        self.__features = []

        # Find `layer-id` text boxes so we have a valid ID **before** using
        # it when setting a shape's `path_id`.
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text
            if notes_text.startswith('.'):
                layer_directive = Parser.layer_directive(notes_text)
                if 'error' in layer_directive:
                    super().error('Slide {}: invalid layer directive: {}'
                                   .format(slide_number, notes_text))
                else:
                    self.set_id(layer_directive.get('id'))
                self.background_for = layer_directive.get('background-for', '')
                self.description = layer_directive.get('description', self.id.capitalize())
                self.models = layer_directive.get('models', '')
                self.outline_feature_id = layer_directive.get('outline')
                self.queryable_nodes = layer_directive.get('queryable-nodes', False)
                self.selectable = self.background_for == '' and not layer_directive.get('not-selectable')
                self.selected = layer_directive.get('selected', False)
                self.zoom = layer_directive.get('zoom', None)
        self.__current_group = []
        # Cannot overlap with slide shape ids...
        self.__next_shape_id = 100001

    @property
    def slide(self):
        return self.__slide

    @property
    def slide_id(self):
        return self.__slide.slide_id

    @property
    def slide_number(self):
        return self.__slide_number

    def feature_id(self, shape_id):
    #==============================
        return '{}#{}'.format(self.slide_id, shape_id)

    def next_feature_id(self):
    #=========================
        id = self.feature_id(self.__next_shape_id)
        self.__next_shape_id += 1
        return id

    def process_initialise(self):
    #============================
        self.__current_group.append('SLIDE')

    def process(self):
    #=================
        self.process_initialise()
        self.process_shape_list(self.slide.shapes, outermost=True)
        self.process_finalise()

    def process_finalise(self):
    #==========================
        pass

    def save(self, filename=None):
    #=============================
        # Override in sub-class
        pass

    def process_group(self, group, properties, *args):
    #=================================================
        self.process_shape_list(group.shapes, *args)

    def process_shape(self, shape, properties, *args):
    #=================================================
        # Override in sub-class
        pass

    def process_shape_list(self, shapes, *args, outermost=False):
    #============================================================
        if not self.selectable:
            return []

        if outermost:
            progress_bar = tqdm(total=len(shapes),
                unit='shp', ncols=40,
                bar_format='{l_bar}{bar}| {n_fmt}/{total_fmt}')

        features = []
        for shape in shapes:
            properties = {
                'shape-name': shape.name,
                'tile-layer': 'features'
                }
            if shape.name.startswith('.'):
                group_name = self.__current_group[-1]  # For error reporting
                properties.update(Parser.shape_markup(shape.name))
                if 'error' in properties:
                    super().error('Shape in slide {}, group {}, has annotation syntax error: {}'
                                  .format(self.__slide_number, group_name, shape.name))
                if 'warning' in properties:
                    super().error('Warning, slide {}, group {}: {}'
                                  .format(self.__slide_number, group_name, properties['warning']))
                for key in ['id', 'path']:
                    if key in properties:
                        if self.mapmaker.duplicate_id(properties[key]):
                           super().error('Shape in slide {}, group {}, has a duplicate id: {}'
                                         .format(self.__slide_number, group_name, shape.name))
            if 'error' in properties:
                pass
            elif 'path' in properties:
                pass
            elif (shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
             or shape.shape_type == MSO_SHAPE_TYPE.FREEFORM
             or isinstance(shape, pptx.shapes.connector.Connector)):
                geometry = self.process_shape(shape, properties, *args)
                feature = Feature(self.feature_id(shape.shape_id), geometry, properties)
                 # Save relationship between id/class and internal feature id
                self.mapmaker.save_feature_id(feature)
                features.append(feature)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self.__current_group.append(properties.get('shape_name', "''"))
                grouped_feature = self.process_group(shape, properties, *args)
                self.__current_group.pop()
                if grouped_feature is not None:
                    self.mapmaker.save_feature_id(grouped_feature)
                    features.append(grouped_feature)
            elif (shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
               or shape.shape_type == MSO_SHAPE_TYPE.PICTURE):
                pass
            else:
                print('"{}" {} not processed...'.format(shape.name, str(shape.shape_type)))
            if outermost:
                progress_bar.update(1)

        if outermost:
            progress_bar.close()
        return features

#===============================================================================

class MapMaker(object):
    def __init__(self, pptx, settings):
        self.__class_to_feature = defaultdict(list)
        self.__id_to_feature = {}
        self.__json_properties = JsonProperties(settings)
        self.__pptx = Presentation(pptx)
        self.__settings = settings
        self.__slides = self.__pptx.slides
        self.__slide_size = [self.__pptx.slide_width, self.__pptx.slide_height]

    def __len__(self):
        return len(self.__slides)

    @property
    def resolved_pathways(self):
        return self.__json_properties.resolved_pathways

    @property
    def settings(self):
        return self.__settings

    @property
    def slide_size(self):
        return self.__slide_size

    def bounds(self):
    #================
        return (0, 0, self.__slide_size[0], self.__slide_size[1])

    def duplicate_id(self, id):
    #==========================
        return self.__id_to_feature.get(id, None) is not None

    def update_properties(self, feature):
    #====================================
        properties = feature.properties.copy()
        properties.update(self.__json_properties.get_properties(properties.get('id'),
                                                                properties.get('class')))
        feature.properties.update(properties)

    def get_slide(self, slide_number):
    #=================================
        slide = self.__slides[slide_number - 1]
        if self.__settings.debug_xml:
            xml = open(os.path.join(self.__settings.output_dir, 'layer{:02d}.xml'.format(slide_number)), 'w')
            xml.write(slide.element.xml)
            xml.close()
        return slide

    def resolve_details(self, layers_dict):
    #======================================
        # Override in sub-class
        pass

    def resolve_pathways(self):
    #==========================
        # Set feature ids of path components
        self.__json_properties.resolve_pathways(self.__id_to_feature, self.__class_to_feature)


    def save_feature_id(self, feature):
    #==================================
        if feature.has_property('id'):
            self.__id_to_feature[feature.property('id')] = feature.feature_id
        if feature.has_property('class'):
            self.__class_to_feature[feature.property('class')].append(feature.feature_id)

#===============================================================================
